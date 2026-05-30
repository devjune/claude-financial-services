#!/usr/bin/env python3
"""
Market Researcher — pptx-author demo.
Builds the sector-primer DECK SHELL that the Market Researcher agent produces,
following the sector-overview skill's 6-step structure. Runs offline via
python-pptx. The live agent fills the [web] placeholders via web search; this
script proves the headless deck-authoring path and shows the target format.

Run:  pip install python-pptx && python3 build_primer_slides.py
      ->  ./out/sector-primer.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt

SECTOR = "Digital Payments / Fintech"
prs = Presentation()  # default template; mount ./templates/firm.pptx for branding

def title_slide(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle
    return s

def bullets_slide(title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    tf = s.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b; p.font.size = Pt(16)
    return s

def table_slide(title, headers, rows):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    nrows, ncols = len(rows) + 1, len(headers)
    tbl = s.shapes.add_table(nrows, ncols, Inches(0.4), Inches(1.6),
                             Inches(9.2), Inches(0.4 * nrows)).table
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = str(val)
    return s

# Step 1 — title / scope
title_slide(f"{SECTOR} — Sector Primer",
            "Illustrative angle: regional digital-wallet & embedded-finance buildout  |  [date]")

# Step 2 — market overview
bullets_slide("Market overview — size & growth", [
    "TAM: [web] $T, 5-yr CAGR [web]%  — cite research firm/methodology",
    "Segmentation: P2P transfer · merchant acquiring · BNPL · embedded finance",
    "Structure: fragmented across regions; top-5 share [web]%",
    "Barriers: licensing/regulation, network effects, merchant distribution",
])

# Step 2b — trends & drivers
bullets_slide("Key trends & drivers", [
    "Tailwinds: cash→digital shift, super-app bundling, real-time payment rails",
    "Headwinds: interchange/regulatory caps, credit-cycle exposure (BNPL)",
    "Disruption: stablecoins, account-to-account (A2A), open banking",
    "Regulation: KYC/AML tightening, data-localization, licensing regimes",
])

# Step 3 — competitive landscape (sector-overview company table)
table_slide("Competitive landscape — key players",
            ["Company", "Revenue", "Growth", "EBITDA margin", "Share", "Differentiator"],
            [["[player A]", "[web]", "[web]", "[web]", "[web]", "merchant network"],
             ["[player B]", "[web]", "[web]", "[web]", "[web]", "super-app reach"],
             ["[player C]", "[web]", "[web]", "[web]", "[web]", "embedded BaaS"],
             ["[player D]", "[web]", "[web]", "[web]", "[web]", "cross-border"]])

# Step 4 — valuation context (comps degrade to [UNSOURCED] without paid MCP)
bullets_slide("Valuation context", [
    "Sector multiples: EV/Revenue [UNSOURCED], EV/EBITDA [UNSOURCED]  — needs CapIQ/FactSet",
    "Premium/discount drivers: growth durability, take-rate, regulatory exposure",
    "Recent M&A transaction multiples: [web / UNSOURCED]",
])

# Step 5 — investment implications / ideas shortlist
bullets_slide("Investment implications — ideas shortlist", [
    "[name 1] — one-line thesis hook",
    "[name 2] — one-line thesis hook",
    "[name 3] — one-line thesis hook",
    "Key debate: take-rate compression vs. volume growth (bull/bear)",
])

os.makedirs("./out", exist_ok=True)
prs.save("./out/sector-primer.pptx")
print(f"wrote ./out/sector-primer.pptx  ({len(prs.slides)} slides)")
